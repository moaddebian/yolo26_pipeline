import sys
sys.stdout = open('extracted_content.txt', 'w', encoding='utf-8')

from pptx import Presentation
prs = Presentation('Synthese_NMS_NMSFree_Stage02_2026.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text)
    print()

print('\n\n========== DOCX CONTENT ==========\n\n')

from docx import Document
doc = Document('NMS-Free dans Yolov26.docx')
for para in doc.paragraphs:
    print(para.text)
for table in doc.tables:
    print('--- TABLE ---')
    for row in table.rows:
        print(' | '.join(cell.text for cell in row.cells))

sys.stdout.close()
print('Done', file=sys.stderr)
